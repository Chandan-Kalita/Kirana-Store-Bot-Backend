from io import BytesIO
from typing import Annotated, Literal, Union

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field, model_validator

_CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}

_NUMBER_FORMATS = {
    "number": "#,##0",
    "currency": '"₹"#,##0',
    "percent": "0%",
}

MAX_TABLE_ROWS = 12
MAX_CHART_CATEGORIES = 8

_HEADER_FILL = RGBColor(0x2F, 0x54, 0x96)


class TitleSlide(BaseModel):
    type: Literal["title"] = "title"
    title: str
    subtitle: str | None = None


class TextSlide(BaseModel):
    type: Literal["text"] = "text"
    title: str
    bullets: list[str]


class TableSlide(BaseModel):
    type: Literal["table"] = "table"
    title: str
    headers: list[str]
    rows: list[list[str]]


class ChartSeries(BaseModel):
    name: str
    values: list[float]


class ChartSlide(BaseModel):
    type: Literal["chart"] = "chart"
    title: str
    chart_type: Literal["bar", "line", "pie"]
    categories: list[str]
    series: list[ChartSeries]
    value_format: Literal["number", "currency", "percent"] = "number"

    @model_validator(mode="after")
    def _check_shape(self) -> "ChartSlide":
        if not self.categories:
            raise ValueError(f"'{self.title}' has no categories.")
        for s in self.series:
            if len(s.values) != len(self.categories):
                raise ValueError(
                    f"'{self.title}': series '{s.name}' has {len(s.values)} "
                    f"values but there are {len(self.categories)} categories "
                    "-- every series must have exactly one value per category."
                )
        if self.chart_type == "pie" and len(self.series) != 1:
            raise ValueError(
                f"'{self.title}': pie charts take exactly one series "
                f"(got {len(self.series)}) -- a pie shows one series split "
                "across categories, not multiple series. Use a bar chart for "
                "multiple series, or split this into separate pie slides."
            )
        if len(self.categories) > MAX_CHART_CATEGORIES:
            raise ValueError(
                f"'{self.title}' has {len(self.categories)} categories, max "
                f"{MAX_CHART_CATEGORIES} -- summarize into fewer buckets "
                "(e.g. top N + 'other') or split across multiple slides."
            )
        return self


# discriminated on "type" -- pydantic renders this as oneOf + discriminator
# in the JSON schema, which Anthropic-style tool calling (incl. DeepSeek's
# Anthropic-compatible endpoint) accepts natively. This is what forces the
# model to produce genuinely structured slide content instead of free text.
SlideSpec = Annotated[
    Union[TitleSlide, TextSlide, TableSlide, ChartSlide], Field(discriminator="type")
]


def _add_title_slide(prs: Presentation, spec: TitleSlide) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = spec.title
    if spec.subtitle:
        slide.placeholders[1].text = spec.subtitle


def _add_text_slide(prs: Presentation, spec: TextSlide) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = spec.title
    body = slide.placeholders[1].text_frame
    body.clear()
    body.word_wrap = True
    # shrinks font-scale instead of clipping if a slide somehow slips past
    # the bullet-count/length guardrail in build_analysis_deck
    body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, bullet in enumerate(spec.bullets):
        paragraph = body.paragraphs[0] if i == 0 else body.add_paragraph()
        paragraph.text = bullet


def _add_table_slide(prs: Presentation, spec: TableSlide) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = spec.title
    rows, cols = len(spec.rows) + 1, len(spec.headers)

    top = Inches(1.5)
    available_height = prs.slide_height - top - Inches(0.3)
    row_height = min(Inches(0.4), available_height // rows)

    table = slide.shapes.add_table(
        rows, cols, Inches(0.5), top, Inches(9), row_height * rows
    ).table
    for r in range(rows):
        table.rows[r].height = row_height

    for c, header in enumerate(spec.headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _HEADER_FILL
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for r, row in enumerate(spec.rows, start=1):
        for c, value in enumerate(row):
            table.cell(r, c).text = value


def _add_chart_slide(prs: Presentation, spec: ChartSlide) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = spec.title
    chart_data = CategoryChartData()
    chart_data.categories = spec.categories
    for series in spec.series:
        chart_data.add_series(series.name, series.values)
    graphic_frame = slide.shapes.add_chart(
        _CHART_TYPES[spec.chart_type],
        Inches(0.5), Inches(1.5), Inches(9), Inches(5),
        chart_data,
    )
    chart = graphic_frame.chart
    chart.font.size = Pt(12)
    number_format = _NUMBER_FORMATS[spec.value_format]

    if spec.chart_type == "pie":
        show_legend = len(spec.categories) > 1
        legend_position = XL_LEGEND_POSITION.RIGHT
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_percentage = spec.value_format == "percent"
        plot.data_labels.show_value = spec.value_format != "percent"
        plot.data_labels.number_format = number_format
        plot.data_labels.number_format_is_linked = False
    else:
        show_legend = len(spec.series) > 1
        legend_position = XL_LEGEND_POSITION.BOTTOM

        category_axis = chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(11)

        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.tick_labels.font.size = Pt(11)
        value_axis.tick_labels.number_format = number_format
        value_axis.tick_labels.number_format_is_linked = False

    # python-pptx's default chart XML ships with a legend already on
    # (position RIGHT) regardless of series count -- explicit off/on both
    # branches, not just the "show" case, or a single-series chart keeps a
    # redundant legend showing only the one series name.
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = legend_position
        # without this the legend can render on top of the plot area
        # instead of beside it, cramming the chart
        chart.legend.include_in_layout = False


_RENDERERS = {
    "title": _add_title_slide,
    "text": _add_text_slide,
    "table": _add_table_slide,
    "chart": _add_chart_slide,
}


def render_deck(title: str, slides: list) -> bytes:
    """Deterministically build a python-pptx presentation from SlideSpec
    objects: one deck title slide up front, then one slide per spec.
    ChartSlide gets a native pptx chart object, not a pasted image.

    If the model's own slides already open with a TitleSlide, that one is
    used instead of auto-generating a second one -- the model's may carry
    a subtitle (e.g. a date range) that the auto-generated one never sets.
    """
    prs = Presentation()
    if slides and isinstance(slides[0], TitleSlide):
        _add_title_slide(prs, slides[0])
        slides = slides[1:]
    else:
        _add_title_slide(prs, TitleSlide(title=title))
    for spec in slides:
        _RENDERERS[spec.type](prs, spec)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
